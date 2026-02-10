import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.colors import ListedColormap, BoundaryNorm
import os

# Fix for Streamlit Cloud (Headless)
plt.switch_backend('Agg')
import traceback

# ==========================================
# CONSTANTES Y CONFIGURACIÓN
# ==========================================
AREA_MINIMA_M2 = 9.0
GRID_SIZE = 1.0

import json

# ==========================================
# CONSTANTES Y CONFIGURACIÓN
# ==========================================
CONFIG_FILE = "config_ai.json"

DEFAULT_SETTINGS = {
    "ai_enabled": True,
    "system_prompt": """Actúa como un Ingeniero Geomensor experto en control de calidad para minería.
Analiza exhaustivamente los datos proporcionados para generar un informe técnico completo.
Estructura tu respuesta en: 
1. Resumen Ejecutivo (Estado general de la poza).
2. Análisis de Zonas Críticas (Identificación de áreas con defectos, sin inventar coordenadas).
3. Recomendaciones Operativas (Acciones correctivas específicas).
IMPORTANTE: NO intentes listar coordenadas específicas si no se te proveen explícitamente (usa referencias a "Tablas adjuntas"). Mantén un tono profesional y directo.""",
    "admin_password": "excon" 
}

def load_settings():
    """Carga configuración de IA desde JSON o devuelve defaults."""
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return {**DEFAULT_SETTINGS, **json.load(f)}
        except:
            return DEFAULT_SETTINGS
    return DEFAULT_SETTINGS

def save_settings(settings):
    """Guarda la configuración actual en JSON."""
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(settings, f, ensure_ascii=False, indent=4)
        return True
    except Exception as e:
        print(f"Error saving settings: {e}")
        return False

# Colores Base (Originales 8 bandas para Estadísticas)
COLORES_FINALES_BASE = [
    '#C00000', # < -3x
    '#FF0000', # -3x a -2x
    '#FFC000', # -2x a -1x
    '#92D050', # -1x a 0
    '#92D050', # 0 a 1x 
    '#00B0F0', # 1x a 2x
    '#0070C0', # 2x a 3x
    '#002060'  # > 3x
]

# Colores Base Simplificados (Para Mapas y Zonas Defectuosas)
COLORES_SEMAFORO = [
    '#FF0000', # Rojo (Bajo Tolerancia)
    '#FFC000', # Amarillo (En Tolerancia pero bajo 0)
    '#00B050'  # Verde (Sobre 0)
]

def get_dynamic_ranges(step):
    """Genera límites y etiquetas basados en el paso de tolerancia (Visual)."""
    limits = [-np.inf, -3*step, -2*step, -1*step, 0, 1*step, 2*step, 3*step, np.inf]
    labels = [
        f'< -{3*step}', f'-{3*step} a -{2*step}', f'-{2*step} a -{step}', 
        f'-{step} a 0', f'0 a {step}', f'{step} a {2*step}', 
        f'{2*step} a {3*step}', f'> {3*step}'
    ]
    return limits, labels

def calcular_rangos(df, rasante=None, step=4.0, dynamic_tol=None):
    """
    Calcula estadísticas de distribución (Tabla Resumen).
    MANTENIENDO LÓGICA V1 (8 Bandas) según solicitud del usuario.
    """
    if df.empty: return pd.DataFrame(), df
    df = df.copy()
    
    # Asegurar existencia de Desv_cm
    if 'desviacion' in df.columns:
        df['Desv_cm'] = df['desviacion']
    elif 'Desv_cm' not in df.columns:
        if rasante is not None:
            df['Desv_cm'] = (df['Cota_Calc'] - rasante) * 100
        else:
            df['Desv_cm'] = 0.0
            
    # Usamos 'step' para la distribución estadística (Visual), no 'dynamic_tol'
    # dynamic_tol se usará SOLO para detectar zonas defectuosas después.
    
    limits, labels_txt = get_dynamic_ranges(step)
            
    rangos = [
        (labels_txt[7], lambda x: x > limits[7], 'Corte Crítico'),
        (labels_txt[6], lambda x: (x > limits[6]) & (x <= limits[7]), 'Corte Alto'),
        (labels_txt[5], lambda x: (x > limits[5]) & (x <= limits[6]), 'Corte Alerta'),
        (labels_txt[4], lambda x: (x >= 0) & (x <= limits[5]), 'OK (Corte)'),
        (labels_txt[3], lambda x: (x >= limits[3]) & (x < 0), 'OK (Relleno)'),
        (labels_txt[2], lambda x: (x >= limits[2]) & (x < limits[3]), 'Relleno Alerta'),
        (labels_txt[1], lambda x: (x >= limits[1]) & (x < limits[2]), 'Relleno Bajo'),
        (labels_txt[0], lambda x: x < limits[1], 'Relleno Crítico')
    ]
    
    res, total = [], len(df)
    for i, (lbl, cond, grp) in enumerate(rangos):
        c = len(df[cond(df['Desv_cm'])])
        # Mapeo de color V1 (8 colores)
        color_idx = 7 - i
        color_hex = COLORES_FINALES_BASE[color_idx]
        
        res.append({
            'Tipo': grp, 'Rango': lbl, 'Puntos': c, 
            'Porcentaje': (c/total)*100 if total > 0 else 0,
            'Color': color_hex
        })
    return pd.DataFrame(res), df

def calculate_dynamic_tolerance(cover_cm):
    """
    Calcula la tolerancia dinámica basada en el espesor (Cover).
    
    Reglas:
    - Espesor > 45 cm: 50% del espesor
    - Espesor 30 a 44.9 cm: 30% del espesor
    - Espesor 20 a 29.9 cm: 10% del espesor
    - Espesor < 20 cm: 0.5 cm (Estricto)
    """
    if pd.isna(cover_cm) or cover_cm <= 0:
        return 0.5 # Valor por defecto seguro (Estricto)
        
    if cover_cm > 45:
        return cover_cm * 0.50
    elif cover_cm >= 30:
        return cover_cm * 0.30
    elif cover_cm >= 20:
        return cover_cm * 0.10
    else:
        return 0.5 # Menor a 20cm, tolerancia cero (usamos 0.5 para estabilidad numérica)

def calcular_rangos_excon(df, cover_cm):
    """
    Calcula estadísticas para CRITERIO EXCON (5 Colores).
    Reglas:
    - Tol_Inf: -15 si Cover > 50, sino -10.
    - Rojo: <= Tol_Inf
    - Amarillo: (Tol_Inf, -4]
    - Verde: (-4, +4]
    - Celeste: (+4, +10]
    - Azul: > +10
    """
    if df.empty: return pd.DataFrame(), df
    df = df.copy()
    
    # Asegurar Desv_cm
    if 'Desv_cm' not in df.columns:
        if 'desviacion' in df.columns: df['Desv_cm'] = df['desviacion']
        else: df['Desv_cm'] = 0.0

    # Determinar Tolerancia Inferior
    tol_inf = -15.0 if cover_cm > 50 else -10.0
    
    # Tolerancia Inferior (as string for labels if needed, but we use strict numeric logic)
    t_inf_val = int(tol_inf)
    
    ranges = [
        # (Label (Range Only), Condition, GroupName (Description), Color)
        (f'<= {t_inf_val}',          lambda x: x <= tol_inf,                  'Critico Bajo',     '#FF0000'), # Rojo
        (f'{t_inf_val} a -4',        lambda x: (x > tol_inf) & (x <= -4),     'Bajo Tolerable',   '#FFC000'), # Amarillo
        ('-4 a 4',                   lambda x: (x > -4) & (x <= 4),           'Conforme',         '#00B050'), # Verde
        ('4 a 10',                   lambda x: (x > 4) & (x <= 10),           'Sobrelevación Leve', '#00B0F0'), # Celeste
        ('> 10',                     lambda x: x > 10,                        'Sobrelevación Crítica', '#002060')  # Azul
    ]
    
    res, total = [], len(df)
    for lbl, cond, grp, color in ranges:
        c = len(df[cond(df['Desv_cm'])])
        res.append({
            'Tipo': grp, 'Rango': lbl, 'Puntos': c, 
            'Porcentaje': (c/total)*100 if total > 0 else 0,
            'Color': color
        })
    
    return pd.DataFrame(res), df

def flood_fill_matrix(matrix):
    """Identifica zonas conectadas en una matriz binaria."""
    rows, cols = matrix.shape
    labeled_matrix = np.zeros_like(matrix, dtype=int)
    current_label = 1
    neighbors = [(-1, -1), (-1, 0), (-1, 1), (0, -1), (0, 1), (1, -1), (1, 0), (1, 1)]
    
    for r in range(rows):
        for c in range(cols):
            if matrix[r, c] == 1 and labeled_matrix[r, c] == 0:
                stack = [(r, c)]
                labeled_matrix[r, c] = current_label
                while stack:
                    curr_r, curr_c = stack.pop()
                    for nr, nc in neighbors:
                        next_r, next_c = curr_r + nr, curr_c + nc
                        if (0 <= next_r < rows and 0 <= next_c < cols and matrix[next_r, next_c] == 1 and labeled_matrix[next_r, next_c] == 0):
                            labeled_matrix[next_r, next_c] = current_label
                            stack.append((next_r, next_c))
                current_label += 1
    return labeled_matrix, current_label - 1

def detectar_zonas(df, col_n, col_e, col_z, tol):
    """Detecta zonas contiguas que exceden la tolerancia. Estrategia: Filtro Puntos -> Grilla."""
    if df.empty: return pd.DataFrame(), 0
    
    # 1. FILTRAR PRIMERO: Solo puntos que son defecto
    # Esto asegura que no se "promedien" defectos con puntos buenos en la misma celda de 1m2
    df_defects = df[df['Desv_cm'] < -tol].copy()
    
    if df_defects.empty:
        return pd.DataFrame(), 0
        
    # Crear grilla solo con puntos defectuosos
    df_defects['GN'] = (df_defects[col_n]//GRID_SIZE)*GRID_SIZE
    df_defects['GE'] = (df_defects[col_e]//GRID_SIZE)*GRID_SIZE
    
    # Agrupar: Z promedio y Desviación promedio (de los malos)
    grid = df_defects.groupby(['GN','GE'])[['Desv_cm', col_z]].mean().reset_index()
    atot = len(grid)*(GRID_SIZE**2) # Area estimada bruta
    
    n_min, e_min = grid['GN'].min(), grid['GE'].min()
    rows = int(grid['GN'].max() - n_min) + 5
    cols = int(grid['GE'].max() - e_min) + 5
    
    if rows > 8000 or cols > 8000: 
        return pd.DataFrame(), atot
    
    mat = np.zeros((rows, cols))
    
    # Map index to grid row and fill matrix
    grid['r_idx'] = ((grid['GN'] - n_min)).astype(int)
    grid['c_idx'] = ((grid['GE'] - e_min)).astype(int)
    
    for _, r in grid.iterrows():
        # Ya filtramos por tolerancia, así que todas las celdas aquí son defecto
        mat[int(r['r_idx']), int(r['c_idx'])] = 1
            
    lbl, num = flood_fill_matrix(mat)
    zonas = []
    
    # Asignar Label a Grid
    grid['Label'] = 0
    for idx, r in grid.iterrows():
        r_i, c_i = int(r['r_idx']), int(r['c_idx'])
        if 0 <= r_i < rows and 0 <= c_i < cols:
            grid.at[idx, 'Label'] = lbl[r_i, c_i]
            
    # Filter only labeled cells (Clusters)
    defect_grid = grid[grid['Label'] > 0]
    
    if not defect_grid.empty:
        zone_grps = defect_grid.groupby('Label')
        
        for label_id, grp in zone_grps:
            area = len(grp) * (GRID_SIZE**2)
            if area >= AREA_MINIMA_M2:
                # Find worst point
                worst_idx = grp['Desv_cm'].idxmin()
                worst_row = grp.loc[worst_idx]
                
                zonas.append({
                    'ID': int(label_id),
                    'Area_Efectiva_m2': area,
                    'Norte': worst_row['GN'],
                    'Este': worst_row['GE'],
                    'Elev_Min': worst_row[col_z],
                    'Desv_Min (cm)': worst_row['Desv_cm']
                })

    return pd.DataFrame(zonas), atot

# ... (Previous imports) ...

def generar_mapa_interactivo(df, zonas_df, col_n, col_e, titulo, tol, criterio="Criterio SQM", cover_cm=0.0):
    """Genera un mapa INTERACTIVO (Plotly) de CALOR (Estilo Técnico XY - Match Excel)."""
    try:
        # Limpieza
        df[col_n] = pd.to_numeric(df[col_n], errors='coerce')
        df[col_e] = pd.to_numeric(df[col_e], errors='coerce')
        df_clean = df.replace([np.inf, -np.inf], np.nan).dropna(subset=[col_n, col_e, 'Desv_cm'])
        
        
        if df_clean.empty: return f"Error: Sin datos válidos."
        
        # Lógica de Color según Criterio
        point_colors = []
        
        if criterio == "Criterio Excon":
            # Tol_Inf depende de Cover
            tol_inf = -15.0 if cover_cm > 50 else -10.0
            
            def get_color_excon(x):
                if x <= tol_inf: return '#FF0000'       # Rojo
                elif x <= -4: return '#FFC000'          # Amarillo
                elif x <= 4: return '#00B050'           # Verde
                elif x <= 10: return '#00B0F0'          # Celeste
                else: return '#002060'                  # Azul
            
            point_colors = df_clean['Desv_cm'].apply(get_color_excon).tolist()
            
        else:
            # Criterio SQM (Original) - Semáforo
            def get_color_semaforo(x):
                if x < -tol: return '#FF0000' # Rojo
                elif x < 0: return '#FFC000'  # Amarillo
                else: return '#00B050'        # Verde
            point_colors = df_clean['Desv_cm'].apply(get_color_semaforo).tolist()

        fig = go.Figure()
        
        # 1. PUNTOS DEL MAPA DE CALOR (Scatter Simple)
        # Usamos Scatter en lugar de Scattergl para asegurar estabilidad visual si el dataset no es masivo.
        fig.add_trace(go.Scatter(
            x=df_clean[col_e],
            y=df_clean[col_n],
            mode='markers',
            marker=dict(
                size=5,
                color=point_colors,
                opacity=0.9,
                line=dict(width=0) # Sin borde para limpieza
            ),
            text=df_clean['Desv_cm'].apply(lambda x: f"Desv: {x:.1f}cm"),
            hoverinfo='text',
            name='Puntos'
        ))
        
        # 2. MARCADORES DE ZONAS DEFECTUOSAS (Negros)
        if not zonas_df.empty:
             if 'Norte' in zonas_df.columns and 'Este' in zonas_df.columns:
                z_n = zonas_df['Norte'].values
                z_e = zonas_df['Este'].values
                z_ids = zonas_df['ID'].values
                z_areas = zonas_df['Area_Efectiva_m2'].values
                z_desv = zonas_df['Desv_Min (cm)'].values if 'Desv_Min (cm)' in zonas_df.columns else [0]*len(z_ids)
                z_elev = zonas_df['Elev_Min'].values if 'Elev_Min' in zonas_df.columns else [0]*len(z_ids)
                
                fig.add_trace(go.Scatter(
                    x=z_e, y=z_n,
                    mode='markers+text',
                    marker=dict(size=12, color='black', symbol='circle', line=dict(color='white', width=1)),
                    text=[str(i) for i in z_ids],
                    textposition='top center',
                    name='Zonas ID',
                    textfont=dict(size=14, color='black', family="Arial Black"),
                    hoverinfo='text',
                    hovertext=[f"ID: {i}<br>Area: {a:.0f} m2<br>Desv Min: {d:.1f} cm<br>Elev Min: {el:.3f} m<br>N: {n:.0f}<br>E: {e:.0f}" 
                               for i, a, d, el, n, e in zip(z_ids, z_areas, z_desv, z_elev, z_n, z_e)]
                ))

        # Configurar Layout Cartesian (Imitando Matplotlib)
        # Fix: Ensure axes are formatted as integers (d) and aspect ratio is 1.
        fig.update_layout(
            # title removed per user request (overlap issues)
            plot_bgcolor='#EBEBEB', # Gris claro de fondo (Estilo Matplotlib)
            xaxis=dict(
                title="Este (X)",
                showgrid=True, gridcolor='white', gridwidth=1,
                zeroline=False,
                scaleanchor="y", scaleratio=1, # Aspect Ratio 1:1
                tickformat="d" # Enteros estrictos
            ),
            yaxis=dict(
                title="Norte (Y)",
                showgrid=True, gridcolor='white', gridwidth=1,
                zeroline=False,
                tickformat="d" # Enteros estrictos (sin comas ni puntos)
            ),
            margin={"r":20,"t":40,"l":20,"b":20},
            height=700, # Un poco más alto para ver mejor
            showlegend=False,
            dragmode='zoom', # Default to zoom
            hovermode='closest'
        )
        return fig
            
    except Exception as e:
        traceback.print_exc()
        return f"Error General Mapa: {str(e)}"

def procesar_turno(df, rasante, tolerancia, col_z, col_n, col_e, step=4.0, cover_cm=0.0, criterio="Criterio SQM"):
    """Procesa un turno completo y retorna resultados."""
    if df.empty:
        return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), 0.0

    # 1. Tolerancia Dinámica (Si hay Cover)
    tol_detect = tolerancia
    # NOTE: topo_dashboard passes calculated tolerance in 'tolerancia' arg already if logic updated there.
    # But if not, we re-calc here? 
    # Current topo_dashboard calls with 'tolerancia=tol_calculated'. So we trust input.
    pass 
    # Logic note: 'calculate_dynamic_tolerance' usage inside here might be redundant if caller handles it.
    # Let's trust the caller provided 'tolerancia' is the correct cut-off.
        
    # 2. Calcular Estadísticas (VISUAL / DISTRIBUCIÓN)
    if criterio == "Criterio Excon":
        tbl_rangos, df_cal = calcular_rangos_excon(df, cover_cm)
    else:
        tbl_rangos, df_cal = calcular_rangos(df, rasante, step=step)
    
    # 3. Detectar Zonas Defectuosas (CRITERIO TÉCNICO / TRAFICO)
    # UPDATED call with col_z
    zonas, area_defectos_bruta = detectar_zonas(df_cal, col_n, col_e, col_z, tol_detect)
    
    # NEW KPI LOGIC (User Request):
    # Incidencia = (Area Zona Defecto / Area TOTAL Trabajada del Turno) * 100
    
    # Calcular Área Total Trabajada (Todo lo levantado/pintado) estimando por grilla
    # Usamos df_cal que ya tiene puntos validos.
    # Grid Size es constante global.
    if not df_cal.empty:
        # Repetimos logica grilla rapida
        gn_all = (df_cal[col_n]//GRID_SIZE).astype(int)
        ge_all = (df_cal[col_e]//GRID_SIZE).astype(int)
        # Unique cells
        unique_cells = len(df_cal.groupby([gn_all, ge_all]).size())
        area_turno_total = unique_cells * (GRID_SIZE**2)
    else:
        area_turno_total = 1.0 # Evitar div/0
        
    # 3. Calcular KPI (Incidencia)
    if not zonas.empty and area_turno_total > 0:
        zonas['KPI Incidencia'] = (zonas['Area_Efectiva_m2'] / area_turno_total)
    elif not zonas.empty:
        zonas['KPI Incidencia'] = 0.0
        
    return tbl_rangos, df_cal, zonas, area_defectos_bruta
    
    # Return tolerance used so dashboard can show it? 
    # Current signature doesn't support returning it, but implementation is enough for now.
    return tbl_rangos, df_cal, zonas, area_tot

def generar_texto_analisis(stats_df, zonas_df, atot, poza):
    """Genera el texto de análisis técnico para el reporte."""
    if stats_df.empty: return "Sin datos."
    
    # Encontrar rango predominante
    pred = stats_df.loc[stats_df['Puntos'].idxmax()]
    
    cant_zonas = len(zonas_df) if not zonas_df.empty else 0
    area_mala = zonas_df['Area_Efectiva_m2'].sum() if not zonas_df.empty else 0
    
    return (f"ANÁLISIS TÉCNICO - {poza}\n\n1. SITUACIÓN GENERAL:\n"
            f"   El rango predominante es '{pred['Rango']}', con un {pred['Porcentaje']:.1f}% de la superficie.\n\n"
            f"2. ÁREAS DEFECTUOSAS:\n   Se detectaron {cant_zonas} zonas críticas (>{AREA_MINIMA_M2}m²). "
            f"La superficie total afectada es de {int(area_mala)} m² sobre un total de {int(atot)} m².\n\n"
            f"3. RECOMENDACIÓN:\n   Se sugiere priorizar las zonas identificadas para trabajos de renivelación.")

# ==========================================
# NUEVA LÓGICA: MAPAS SATELITALES (MATPLOTLIB + CONTEXTILY)
# ==========================================
import contextily as ctx
try:
    import pyproj
except ImportError:
    pass

def generar_mapa_satelital_interactivo(df, zonas_df, col_n, col_e, titulo, tol):
    """Genera un mapa INTERACTIVO (Plotly) con fondo Satelital (Esri)."""
    try:
        # Limpieza y Conversión Numérica
        df[col_n] = pd.to_numeric(df[col_n], errors='coerce')
        df[col_e] = pd.to_numeric(df[col_e], errors='coerce')
        
        df_clean = df.replace([np.inf, -np.inf], np.nan).dropna(subset=[col_n, col_e, 'Desv_cm'])
        if df_clean.empty: 
            return f"Error: Sin datos válidos. Revise columnas {col_n}/{col_e} o valores vacíos."
        
        # CONVERSIÓN DE COORDENADAS (UTM 19S -> WGS84 Lat/Lon)
        # Requerido para mapas web (Plotly Mapbox)
        transformer = None
        try:
            transformer = pyproj.Transformer.from_crs("EPSG:32719", "EPSG:4326", always_xy=True)
            # Transform expects (x, y) -> returns (lon, lat)
            lons, lats = transformer.transform(df_clean[col_e].values, df_clean[col_n].values)
            df_clean['lon'] = lons
            df_clean['lat'] = lats
        except Exception as e:
            return f"Error Proyección: {str(e)}"

        # Configurar colores y rangos (igual que antes)
        limits, labels = get_dynamic_ranges(tol)
        
        fig = go.Figure()
        
        # --- ZONAS DEFECTUOSAS MARCADORES ---
        if not zonas_df.empty:
             if 'Norte' in zonas_df.columns and 'Este' in zonas_df.columns:
                z_n = zonas_df['Norte'].values
                z_e = zonas_df['Este'].values
                z_ids = zonas_df['ID'].values
                
                if transformer:
                    lat_z, lon_z = transformer.transform(z_e, z_n)
                else:
                    lat_z, lon_z = z_n, z_e
                    
                fig.add_trace(go.Scattermapbox(
                    lat=lat_z, lon=lon_z,
                    mode='markers+text',
                    marker=dict(size=14, color='black', symbol='circle'),
                    text=[str(i) for i in z_ids],
                    textposition='top center',
                    name='Zonas ID',
                    textfont=dict(size=14, color='black', family="Arial Black"),
                    hoverinfo='text',
                    hovertext=[f"ID: {i}<br>Area: {a:.0f}m2<br>Desv Min: {d:.1f}cm" 
                               for i, a, d in zip(z_ids, zonas_df['Area_Efectiva_m2'], zonas_df['Desv_Min (cm)'])]
                ))


        # Añadir Puntos
        fig.add_trace(go.Scattermapbox(
            lat=df_clean['lat'],
            lon=df_clean['lon'],
            mode='markers',
            marker=go.scattermapbox.Marker(
                size=8, # Tamaño pixelado visible
                color=df_clean['Desv_cm'],
                colorscale=[
                    [0.0, '#C00000'],    # < -3x
                    [0.125, '#FF0000'],  # -3x a -2x
                    [0.25, '#FFC000'],   # -2x a -1x
                    [0.375, '#92D050'],  # -1x a 0
                    [0.5, '#92D050'],    # 0 a 1x 
                    [0.625, '#00B0F0'],  # 1x a 2x
                    [0.75, '#0070C0'],   # 2x a 3x
                    [1.0, '#002060']     # > 3x
                ],
                cmin=-3*tol, cmax=3*tol, # Fijar min/max para estabilidad visual
                opacity=0.8,
            ),
            text=df_clean['Desv_cm'].apply(lambda x: f"Desv: {x:.1f}cm"),
            hoverinfo='text'
        ))

        # Configurar Layout Maps (Esri Satellite)
        fig.update_layout(
            title=dict(text=f"Satelital Interactivo - {titulo}", y=0.98),
            mapbox=dict(
                style="white-bg", # Estilo base vacío para poner capas encima
                layers=[
                    {
                        "below": 'traces',
                        "sourcetype": "raster",
                        "sourceattribution": "Esri World Imagery",
                        "source": [
                            "https://server.arcgisonline.com/ArcGIS/rest/services/World_Imagery/MapServer/tile/{z}/{y}/{x}"
                        ]
                    }
                ],
                center=dict(lat=df_clean['lat'].mean(), lon=df_clean['lon'].mean()),
                zoom=16 # Zoom inicial cercano
            ),
            margin={"r":0,"t":40,"l":0,"b":0},
            height=600
        )
        return fig
            
    except Exception as e:
        traceback.print_exc()
        return f"Error General Satélite: {str(e)}"

# ==========================================
# LÓGICA DE VISUALIZACIÓN (ANTIGUA)
# ==========================================
from matplotlib.colors import ListedColormap, BoundaryNorm

def generar_mapa_matplotlib(df, zonas, col_n, col_e, titulo, tol, criterio="Criterio Excon"):
    """
    Genera un objeto Figura de Matplotlib con el mapa de calor.
    Retorna (fig, ax) para ser usado en Streamlit o Reportes.
    """
    try:
        # Limpieza robusta y Coerción Numérica
        df[col_n] = pd.to_numeric(df[col_n], errors='coerce')
        df[col_e] = pd.to_numeric(df[col_e], errors='coerce')

        df_clean = df.replace([np.inf, -np.inf], np.nan).dropna(subset=[col_n, col_e, 'Desv_cm'])
        if df_clean.empty:
            raise ValueError(f"Sin datos válidos para graficar. Revisar columnas {col_n}, {col_e} o Desv_cm.")

        x_min, x_max = df_clean[col_e].min(), df_clean[col_e].max()
        y_min, y_max = df_clean[col_n].min(), df_clean[col_n].max()
        
        # Validación de rangos
        if pd.isna(x_min) or pd.isna(x_max) or pd.isna(y_min) or pd.isna(y_max):
             raise ValueError(f"Rango de coordenadas inválido (NaN). Min/Max E: {x_min}/{x_max}, N: {y_min}/{y_max}")

        # Margen dinámico (5% del rango o min 5m)
        dx, dy = x_max - x_min, y_max - y_min
        margin = max(5, max(dx, dy) * 0.05)

        # Configurar Figura
        plt.style.use('ggplot')
        fig, ax = plt.subplots(figsize=(10, 8), dpi=100) 
        ax.set_facecolor('white') # Explicit White Background
        
        # --- COLOR LOGIC BASED ON CRITERION ---
        if criterio == "Criterio Excon":
            # Excon Logic: 
            # Red: <= -Tol (Crítico Bajo)
            # Yellow: > -Tol and <= -4 (Bajo Tolerable)
            # Green: > -4 and <= 4 (Conforme)
            # Light Blue: > 4 and <= 10 (Sobrelev. Leve)
            # Dark Blue: > 10 (Sobrelev. Crítica)
            
            if tol < 4: tol = 15.0 # Logic implies Tol is the outer limit.

            bins = [-np.inf, -tol, -4, 4, 10, np.inf]
            color_list = [
                '#FF0000', # Red 
                '#FFC000', # Yellow
                '#00B050', # Green
                '#00B0F0', # Light Blue
                '#002060'  # Dark Blue
            ]
            
            legend_elements = [
                patches.Patch(facecolor='#FF0000', label=f'Crítico Bajo (<= -{int(tol)})'),
                patches.Patch(facecolor='#FFC000', label=f'Bajo Tolerable (-{int(tol)} a -4)'),
                patches.Patch(facecolor='#00B050', label=f'Conforme (-4 a +4)'),
                patches.Patch(facecolor='#00B0F0', label=f'Sobrelev. Leve (+4 a +10)'),
                patches.Patch(facecolor='#002060', label=f'Sobrelev. Crítica (> +10)')
            ]

        else: # Criterio SQM (Standard Tolerance)
            # Red: Outside Tolerance (+/- Tol)
            # Green: Inside Tolerance
            
            # We can use a simpler 3-bin or continuous approach, but let's map to discrete for consistency
            # <-Tol (Red), -Tol to Tol (Green), >Tol (Red)
            # Or maybe detailed? Usually SQM is strict Pass/Fail.
            # Let's use a 3-color scheme for visuals: Red (Low), Green (OK), Red (High)
            
            bins = [-np.inf, -tol, tol, np.inf]
            color_list = [
                '#FF0000', # Red (Low)
                '#00B050', # Green (OK)
                '#FF0000'  # Red (High)
            ]

             # Alternative: Blue for High to distinguish?
            color_list = [
                '#FF0000', # Red (Low)
                '#00B050', # Green (OK)
                '#0000FF'  # Blue (High) - visual distinction
            ]

            legend_elements = [
                patches.Patch(facecolor='#FF0000', label=f'Bajo Tolerancia (< -{tol})'),
                patches.Patch(facecolor='#00B050', label=f'Conforme (-{tol} a +{tol})'),
                patches.Patch(facecolor='#0000FF', label=f'Sobre Tolerancia (> +{tol})')
            ]

        
        cmap_custom = ListedColormap(color_list)
        norm = BoundaryNorm(bins, cmap_custom.N)
        
        # Scatter Plot
        sc = ax.scatter(
            df_clean[col_e], df_clean[col_n],
            c=df_clean['Desv_cm'],
            cmap=cmap_custom, norm=norm,
            s=20, marker='o', alpha=0.9, edgecolors='none', zorder=10
        )

        # Custom Legend (Outside Right)
        # Using bbox_to_anchor to push it further right
        leg = ax.legend(
            handles=legend_elements, 
            title="Clasificación (cm)", 
            bbox_to_anchor=(1.02, 1), 
            loc='upper left', 
            borderaxespad=0.,
            frameon=True, 
            facecolor='white', 
            framealpha=1.0
        )
        leg.get_title().set_fontsize('8') 
        for text in leg.get_texts():
            text.set_fontsize('7')
            
        # Fix Coordinates Display
        ax.ticklabel_format(useOffset=False, style='plain')
        ax.set_xlabel("Este (X)", fontsize=8)
        ax.set_ylabel("Norte (Y)", fontsize=8)
        
        # ROTATE LABELS to avoid overlapping ("amontonadas")
        ax.tick_params(axis='x', rotation=45, labelsize=7)
        ax.tick_params(axis='y', labelsize=7)
        
        # Ensure layout accommodates external legend
        # Adjust margins explicitly
        plt.subplots_adjust(right=0.85) # Leave space for legend
        
        # grid blanco fino
        ax.grid(True, color='white', linestyle='-', linewidth=0.5, alpha=0.8)
        
        # Titulo y Ejes
        ax.set_title(f"{titulo}", fontsize=14, fontweight='bold', pad=15)
        ax.set_aspect('equal')
        ax.set_xlim(x_min - margin, x_max + margin)
        ax.set_ylim(y_min - margin, y_max + margin)

        if not zonas.empty:
            # --- MARCADORES DE PUNTOS CRÍTICOS (ID) ---
            # Check if we have point coordinates (New Logic)
            if 'Norte' in zonas.columns and 'Este' in zonas.columns:
                 # Plot markers (Black dots)
                 ax.scatter(
                     zonas['Este'], zonas['Norte'],
                     c='black', s=10, marker='o', edgecolors='white', linewidth=1, zorder=25, label='Zona Defectuosa'
                 )
                 # Add ID labels
                 for _, z in zonas.iterrows():
                     ax.text(
                         z['Este'], z['Norte'], str(int(z['ID'])),
                         color='black', fontsize=9, fontweight='bold', ha='center', va='bottom', zorder=30
                     )
            # Fallback for old data (Rectangles)
            elif 'E_Min' in zonas.columns:
                for _, z in zonas.iterrows():
                    if pd.isna(z['E_Min']) or pd.isna(z['N_Min']): continue
                    width = z['E_Max'] - z['E_Min']
                    height = z['N_Max'] - z['N_Min']
                    rect = patches.Rectangle(
                        (z['E_Min'], z['N_Min']), width, height,
                        linewidth=2.0, edgecolor='#FF0000', facecolor='none', zorder=20
                    )
                    ax.add_patch(rect)

        return fig
        
    except Exception:
        traceback.print_exc()
        return None

# ==========================================
# NUEVA LÓGICA: MAPAS INTERACTIVOS (PLOTLY)
# ==========================================
import plotly.graph_objects as go

# Legacy function removed. Now standardized on generar_mapa_interactivo (formerly satelital).

# Legacy function removed. Now standardized on generar_mapa_interactivo (formerly satelital).

# ==========================================
# POWERPOINT GENERATION (Restored)
# ==========================================
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import copy
import os
from datetime import datetime
import matplotlib.pyplot as plt
import matplotlib.patches as patches

def duplicate_slide(pres, index):
    try:
        template = pres.slides[index]
    except IndexError:
        return pres.slides.add_slide(pres.slide_layouts[6])

    # Create empty slide using Blank layout (usually index 6 or find by name)
    try:
        empty_layout = pres.slide_layouts[6] 
    except:
        empty_layout = pres.slide_layouts[-1]
        
    slide = pres.slides.add_slide(empty_layout)
    
    # Copy shapes
    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return slide

def _create_clean_slide(prs, template_index, poza_id, turno, subtitle):
    """Helper to clone and aggressively clean a slide, preserving Title."""
    try:
        slide = duplicate_slide(prs, template_index)
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
    # Aggressive Clean: Remove known junk AND placeholders
    junk_ids = [20, 24, 26, 28, 30, 36, 40, 25, 35, 42, 27, 31, 29, 22, 12, 13, 16, 4, 39, 23, 2, 6, 9, 11, 13, 10, 19] 
    
    shapes_to_delete = []
    
    # Iterate copy
    for shp in slide.shapes:
        # Preserve Title (ID 3)
        if shp.shape_id == 3:
            continue
            
        # Delete junk by ID
        if shp.shape_id in junk_ids:
            shapes_to_delete.append(shp)
            continue
            
        # Delete small textboxes (Numbers, labels)
        if shp.has_text_frame:
             txt = shp.text_frame.text.strip()
             if len(txt) < 3 and txt.isdigit(): 
                 shapes_to_delete.append(shp)
                 continue
             # Remove specific placeholder texts if found
             if "registro" in txt.lower() or "cumple" in txt.lower():
                 shapes_to_delete.append(shp)
                 continue

    for shp in shapes_to_delete:
        sp = shp.element
        if sp.getparent() is not None:
             sp.getparent().remove(sp)
             
    # Set Title
    title_set = False
    for shp in slide.shapes:
        if shp.shape_id == 3:
             if shp.has_text_frame:
                shp.text_frame.text = f"Reporte Control: {poza_id} - Turno {turno} - {subtitle}"
                for p in shp.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255) # White
                title_set = True
                
    if not title_set:
        # Fallback Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.8))
        tb.text_frame.text = f"Reporte Control: {poza_id} - Turno {turno} - {subtitle}"
        p = tb.text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0,0,0) # Fallback black if no blue bar
        
    return slide

def generar_pptx_report(global_res_dict, template_path="Regis_GPS_25_Enero_Turno_B_Cosecha.pptx"):
    output = io.BytesIO()
    try:
        if os.path.exists(template_path):
            prs = Presentation(template_path)
            # USE SLIDE 3 (Index 2) AS TEMPLATE (Has "Registro acumulado" vs "Registro turno día")
            TEMPLATE_INDEX = 2 
        else:
            prs = Presentation() 
            TEMPLATE_INDEX = 0
            prs.slides.add_slide(prs.slide_layouts[0])
            
    except Exception as e:
        print(f"Error loading template: {e}")
        prs = Presentation()
        TEMPLATE_INDEX = 0

    # Iterate Pozas and Turns
    original_slide_count = len(prs.slides)
    
    for poza_id, poza_data in global_res_dict.items():
        # Only process Turns A and B (No General / Acumulado)
        current_date_str = datetime.now().strftime("%d/%m/%Y")
        
        for t in ['A', 'B']:
            if t not in poza_data or poza_data[t]['vacio']: continue
            data = poza_data[t]
            
            # Extract Analysis Parts
            part_general = data.get('texto_analisis', "Sin comentarios.")
            part_critical = data.get('texto_analisis_critico', "Sin comentarios críticos.")
            
            # Fallback for old data structure (Splitting)
            if part_critical == "Sin comentarios críticos." and "2. ÁREAS DEFECTUOSAS:" in part_general:
                 try:
                     parts = part_general.split("2. ÁREAS DEFECTUOSAS:")
                     part_general = parts[0].replace("1. SITUACIÓN GENERAL:", "").replace(f"ANÁLISIS TÉCNICO - {poza_id}", "").strip()
                     part_critical = "2. ÁREAS DEFECTUOSAS:" + parts[1]
                 except: pass

            # ====================================================
            # SLIDE 1: RESUMEN EJECUTIVO
            # Layout: Title -> Text (Top) -> Stats Table (Left Bottom) -> Chart (Right Bottom)
            # ====================================================
            slide1 = _create_clean_slide(prs, TEMPLATE_INDEX, poza_id, t, f"Resumen Ejecutivo - {current_date_str}")
            
            # 1. General Analysis Text (Top)
            tb_gen = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(2.5))
            tf = tb_gen.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = "SITUACIÓN GENERAL:"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0,0,0)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(3)
            
            p2 = tf.add_paragraph()
            p2.text = part_general
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(0,0,0)
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = 1.0 # Single spacing
            p2.space_after = Pt(6)

            # 2. Stats Table (Left Bottom)
            # Add Title for Table
            tb_tbl_ti = slide1.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(4.5), Inches(0.3))
            p_ti = tb_tbl_ti.text_frame.add_paragraph()
            p_ti.text = "Distribución de Puntos por Rango"
            p_ti.font.size = Pt(10)
            p_ti.font.bold = True
            p_ti.alignment = PP_ALIGN.LEFT
            
            if 'tbl' in data and not data['tbl'].empty:
                df_tbl = data['tbl'].copy()
                cols_order = ['Tipo', 'Rango', 'Puntos', 'Porcentaje']
                # Ensure columns exist
                cols_exist = [c for c in cols_order if c in df_tbl.columns]
                df_tbl = df_tbl[cols_exist]
                
                rows_t, cols_t = df_tbl.shape
                if rows_t > 0:
                    table_shape = slide1.shapes.add_table(rows_t+1, cols_t, Inches(1.5), Inches(4.5), Inches(4.5), Inches(2.0))
                    table = table_shape.table
                    for i, col_name in enumerate(df_tbl.columns):
                        cell = table.cell(0, i)
                        cell.text = str(col_name)
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(9)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                        p.font.color.rgb = RGBColor(0,0,0)
                    for r_idx, row in df_tbl.iterrows():
                        for c_idx, val in enumerate(row):
                             cell = table.cell(r_idx+1, c_idx)
                             cell.text = f"{val:.1f}%" if 'Porcentaje' in df_tbl.columns[c_idx] else str(val)
                             p = cell.text_frame.paragraphs[0]
                             p.font.size = Pt(9)
                             p.alignment = PP_ALIGN.CENTER
                             p.font.color.rgb = RGBColor(0,0,0)

            # 3. Bar Chart logic
            if 'tbl' in data:
                 try:
                     fig_bar, ax = plt.subplots(figsize=(5, 3))
                     colors = data['tbl']['Color'].tolist() if 'Color' in data['tbl'] else ['blue']*len(data['tbl'])
                     ax.bar(data['tbl']['Rango'], data['tbl']['Puntos'], color=colors)
                     ax.set_title("Distribución de Puntos por Rango", fontsize=10, pad=10)
                     ax.tick_params(axis='x', rotation=45, labelsize=8)
                     ax.tick_params(axis='y', labelsize=8)
                     ax.spines['top'].set_visible(False)
                     ax.spines['right'].set_visible(False)
                     
                     img_bar = io.BytesIO()
                     fig_bar.savefig(img_bar, format='png', bbox_inches='tight', dpi=120)
                     img_bar.seek(0)
                     plt.close(fig_bar)
                     slide1.shapes.add_picture(img_bar, Inches(7.0), Inches(4.2), width=Inches(4.5))
                 except: pass

            # ====================================================
            # SLIDE 2: DETALLE DE DEFECTOS
            # ====================================================
            slide2 = _create_clean_slide(prs, TEMPLATE_INDEX, poza_id, t, f"Detalle de Defectos - {current_date_str}")
            
            # 1. Critical Analysis Text (Top Left)
            # GENERATE STANDARDIZED STATS TEXT
            # Fetch Stats
            area_tot = data.get('atot', 0)
            if area_tot == 0 and 'Stats' in poza_data:
                area_tot = poza_data['Stats'].get('Area_Total_m2', 0)
            
            stats_paragraph = f"La superficie total trabajada en este turno corresponde a {area_tot:,.0f} m²."
            
            has_zones = ('zonas' in data and not data['zonas'].empty)
            if has_zones:
                 area_def = data['zonas']['Area_Efectiva_m2'].sum() if 'Area_Efectiva_m2' in data['zonas'].columns else 0
                 incidencia = (area_def / area_tot * 100.0) if area_tot > 0 else 0
                 stats_paragraph += f" Se detectaron {len(data['zonas'])} zonas de puntos bajos que suman {area_def:,.0f} m², representando una incidencia del {incidencia:.1f}% respecto a la superficie total."
            else:
                 stats_paragraph += " No se detectaron zonas críticas que superen el criterio de área mínima, logrando un 100% de cumplimiento en la superficie analizada."

            # Merge with existing analysis if any, or replace empty "Sine comentarios"
            if part_critical == "Sin comentarios críticos." or part_critical == "No hay zonas críticas para analizar.":
                full_critical_text = stats_paragraph
            else:
                full_critical_text = f"{stats_paragraph}\n\nDetalle Adicional: {part_critical}"

            tb_crit = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(2.5))
            tf = tb_crit.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = full_critical_text
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0,0,0)
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.0
            p.space_after = Pt(6)

            # 2. Critical Zones Table (Left Bottom)
            # Title only if zones exist
            if has_zones:
                tb_zt = slide2.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(4.0), Inches(0.5))
                tb_zt.text = "Puntos Bajos Detectados"
                tb_zt.text_frame.paragraphs[0].font.size = Pt(10)
                tb_zt.text_frame.paragraphs[0].font.bold = True
                
                df_z = data['zonas'].copy()
                cols_target = ['ID', 'Area_Efectiva_m2', 'Norte', 'Este', 'Elev_Min', 'Desv_Min (cm)']
                cols_show = [c for c in cols_target if c in df_z.columns]
                df_z_show = df_z[cols_show].head(12)
                rows_z, cols_z = df_z_show.shape
                
                if rows_z > 0:
                    table_z_shape = slide2.shapes.add_table(rows_z+1, cols_z, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.5))
                    table_z = table_z_shape.table
                    # FORCE HEADER ROW HEIGHT
                    table_z.rows[0].height = Inches(0.20)
                    
                    # Header Logic
                    for i, c_name in enumerate(cols_show):
                        cell = table_z.cell(0, i)
                        cell.text = c_name
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(8)
                        p.bold = True
                        p.alignment = PP_ALIGN.CENTER
                    # Data Row Logic + HEIGHT REDUCTION
                    for r_idx, row in df_z_show.iterrows():
                        # Set Row Height (Half Size ~ 0.2 inch)
                        table_z.rows[r_idx+1].height = Inches(0.20)
                        
                        for c_idx, col_name in enumerate(cols_show):
                             val = row[col_name]
                             cell = table_z.cell(r_idx+1, c_idx)
                             # Value formatting...
                             if isinstance(val, float):
                                 if 'ID' in col_name: fmt = "{:.0f}"
                                 elif 'Elev' in col_name: fmt = "{:.3f}"
                                 elif 'Desv' in col_name: fmt = "{:.2f}"
                                 elif 'Area' in col_name or 'Norte' in col_name or 'Este' in col_name: fmt = "{:.0f}"
                                 else: fmt = "{:.2f}"
                                 cell.text = fmt.format(val)
                             else:
                                 cell.text = str(val)
                                 
                             p = cell.text_frame.paragraphs[0]
                             p.font.size = Pt(8)
                             p.alignment = PP_ALIGN.CENTER
                             p.font.color.rgb = RGBColor(0,0,0)
                             # Reduce margins for compact look
                             cell.margin_top = Pt(1)
                             cell.margin_bottom = Pt(1)
            else:
                 pass

            # 3. Heatmap (Right Bottom)
            try:
                # FORCE BACKEND SWITCH
                plt.switch_backend('Agg') 
                
                # Robust column detection (Created in dashboard)
                cn_map = 'Norte' if 'Norte' in data['df'].columns else ('N' if 'N' in data['df'].columns else data['df'].columns[1])
                ce_map = 'Este' if 'Este' in data['df'].columns else ('E' if 'E' in data['df'].columns else data['df'].columns[0])

                # DETECT CRITERION FROM CONFIG
                crit_map = poza_data.get('Config', {}).get('Criterio', 'Criterio Excon')

                fig = generar_mapa_matplotlib(
                     data['df'], data['zonas'], 
                     col_n=cn_map, 
                     col_e=ce_map,
                     titulo=f"Mapa - Turno {t} ({crit_map})",
                     tol=poza_data['Config']['Tol'],
                     criterio=crit_map
                )
                
                # Add Frame/Border
                fig.patch.set_linewidth(1)
                fig.patch.set_edgecolor('black')
                
                # Check for Figure vs Error String vs None
                if fig and hasattr(fig, 'savefig'):
                    fig.patch.set_facecolor('white')
                    img_stream = io.BytesIO()
                    # High quality (300 DPI) and Tight Bounding Box
                    fig.savefig(img_stream, format='png', bbox_inches='tight', dpi=300, facecolor='white', edgecolor='black')
                    img_stream.seek(0)
                    plt.close(fig)
                    
                    slide2.shapes.add_picture(img_stream, Inches(7.5), Inches(0.9), height=Inches(5.9))
                else:
                     # Handle error return (string or None)
                     err_msg = fig if isinstance(fig, str) else "Generador retornó None"
                     tb = slide2.shapes.add_textbox(Inches(7.0), Inches(4.5), Inches(5.0), Inches(1.0))
                     tb.text = f"(Error Gen Mapa: {err_msg})"
                     p = tb.text_frame.paragraphs[0]
                     p.font.size = Pt(10)
                     p.font.color.rgb = RGBColor(255,0,0)

            except Exception as e:
                print(f"Map Error: {e}")
                tb = slide2.shapes.add_textbox(Inches(7.0), Inches(4.5), Inches(5.0), Inches(2.0))
                tb.text = f"Error Crítico Mapa: {str(e)}"
            

    # DELETE ORIGINAL TEMPLATE SLIDES (KEEP COVER - Index 0)
    if original_slide_count > 1:
        for _ in range(original_slide_count - 1):
            if len(prs.slides) > 1:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                if len(slides) > 1:
                     xml_slides.remove(slides[1]) # Remove index 1 continuously

    prs.save(output)
    return output.getvalue()

